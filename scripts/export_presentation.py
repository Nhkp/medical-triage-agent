from __future__ import annotations

import argparse
import html.parser
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any


@dataclass
class Slide:
    title: str = ""
    eyebrow: str = ""
    subtitle: str = ""
    bullets: list[str] = field(default_factory=list)
    metrics: list[str] = field(default_factory=list)


class DeckParser(html.parser.HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.slides: list[Slide] = []
        self._current: Slide | None = None
        self._capture: str | None = None
        self._buffer: list[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        attrs_dict = dict(attrs)
        classes = set((attrs_dict.get("class") or "").split())
        if tag == "section" and "slide" in classes:
            self._current = Slide(title=attrs_dict.get("data-slide-title") or "")
        elif self._current and tag in {"h1", "h2", "h3", "p", "li", "strong", "span"}:
            self._capture = tag
            self._buffer = []
        elif tag == "aside":
            self._skip_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag == "aside" and self._skip_depth:
            self._skip_depth -= 1
            return
        if tag == "section" and self._current:
            self.slides.append(self._current)
            self._current = None
            return
        if self._current and self._capture == tag:
            text = _clean("".join(self._buffer))
            if text:
                self._assign_text(tag, text)
            self._capture = None
            self._buffer = []

    def handle_data(self, data: str) -> None:
        if self._capture and not self._skip_depth:
            self._buffer.append(data)

    def _assign_text(self, tag: str, text: str) -> None:
        if not self._current:
            return
        if tag in {"h1", "h2"}:
            self._current.title = text
        elif tag in {"h3", "li"}:
            self._current.bullets.append(text)
        elif tag == "p" and not self._current.eyebrow:
            self._current.eyebrow = text
        elif tag == "p" and not self._current.subtitle:
            self._current.subtitle = text
        elif tag in {"strong", "span"}:
            self._current.metrics.append(text)


def main() -> int:
    args = _parse_args()
    slides = parse_deck(Path(args.input))
    if args.dry_run:
        print(f"slides={len(slides)}")
        for index, slide in enumerate(slides, start=1):
            print(f"{index}. {slide.title}")
        return 0
    export_pptx(slides, Path(args.output))
    print(f"wrote {args.output}")
    return 0


def parse_deck(path: Path) -> list[Slide]:
    parser = DeckParser()
    parser.feed(path.read_text(encoding="utf-8"))
    if not parser.slides:
        raise ValueError(f"no slides found in {path}")
    return parser.slides


def export_pptx(slides: list[Slide], output_path: Path) -> None:
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
        from pptx.util import Inches, Pt  # type: ignore[import-untyped]
    except ImportError as exc:
        raise RuntimeError(
            "Install presentation extras with `uv sync --extra presentation`."
        ) from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for index, slide_data in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        _background(slide, prs, RGBColor(7, 19, 18))
        _accent_bar(slide, prs, RGBColor(64, 214, 176), index)
        _textbox(
            slide,
            slide_data.eyebrow,
            0.75,
            0.55,
            9.8,
            0.35,
            Pt(10),
            RGBColor(64, 214, 176),
            bold=True,
        )
        _textbox(
            slide,
            slide_data.title,
            0.72,
            1.05,
            10.8,
            1.6,
            Pt(34),
            RGBColor(238, 247, 247),
            bold=True,
        )
        if slide_data.subtitle:
            _textbox(
                slide, slide_data.subtitle, 0.78, 2.65, 10.8, 0.85, Pt(16), RGBColor(159, 184, 183)
            )
        _content(slide, slide_data, RGBColor(238, 247, 247), RGBColor(159, 184, 183))
        _textbox(slide, f"{index:02d}", 12.15, 6.8, 0.7, 0.25, Pt(9), RGBColor(159, 184, 183))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def _content(slide: Any, slide_data: Slide, ink: Any, muted: Any) -> None:
    top = 3.65
    if slide_data.metrics:
        values = _metric_pairs(slide_data.metrics)
        left = 0.78
        for value, label in values[:4]:
            _textbox(slide, value, left, top, 2.6, 0.55, _pt(29), ink, bold=True)
            _textbox(slide, label, left, top + 0.58, 2.6, 0.4, _pt(11), muted)
            left += 2.95
        top += 1.35
    if slide_data.bullets:
        body = "\n".join(f"- {bullet}" for bullet in slide_data.bullets[:7])
        _textbox(slide, body, 0.82, top, 11.25, 2.7, _pt(16), ink)


def _metric_pairs(items: list[str]) -> list[tuple[str, str]]:
    pairs: list[tuple[str, str]] = []
    pending: str | None = None
    for item in items:
        if pending is None:
            pending = item
        else:
            pairs.append((pending, item))
            pending = None
    return pairs if pairs else [(item, "") for item in items]


def _background(slide: Any, prs: Any, color: Any) -> None:
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore[import-untyped]

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _accent_bar(slide: Any, prs: Any, color: Any, index: int) -> None:
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore[import-untyped]
    from pptx.util import Inches  # type: ignore[import-untyped]

    width = Inches(0.08)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, prs.slide_width - width, 0, width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if index % 3 == 0:
        shape.fill.transparency = 18


def _textbox(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: Any,
    color: Any,
    *,
    bold: bool = False,
) -> None:
    from pptx.enum.text import PP_ALIGN  # type: ignore[import-untyped]
    from pptx.util import Inches  # type: ignore[import-untyped]

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color


def _pt(value: int) -> Any:
    from pptx.util import Pt  # type: ignore[import-untyped]

    return Pt(value)


def _clean(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Export the HTML project presentation to PPTX")
    parser.add_argument("--input", default="presentations/chsa-current-state/index.html")
    parser.add_argument("--output", default="dist/presentations/chsa-current-state.pptx")
    parser.add_argument("--dry-run", action="store_true")
    return parser.parse_args()


if __name__ == "__main__":
    raise SystemExit(main())
